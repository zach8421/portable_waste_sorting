"""Build the final-project presentation deck.

Generates ``Portable_Waste_Sorting.pptx`` at the repo root from the content in
this docs/ folder. Re-run after editing the SLIDES data below:

    python3 docs/build_pptx.py

Covers the five required summary points: information story, existing structure
+ FAIR assessment, how we decided to improve, the new structure, and how
quality is identified and addressed.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(REPO, "Portable_Waste_Sorting.pptx")

GREEN = RGBColor(0x2E, 0x7D, 0x32)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)

# (title, subtitle-or-None, [ (text, level), ... ])
SLIDES = [
    ("Portable Waste Sorting Information for Seattle",
     "IMT 542 Final Project  ·  Information Architect  ·  github.com/zach8421/portable_waste_sorting",
     None),

    ("1 · The Information Story", "Who needs this, and why it matters", [
        ("User: a Seattle resident at the bin asking \"which bin does this go in?\" — plus developers, educators, and civic-tech builders who reuse the guidance.", 0),
        ("Problem: SPU's official answers are locked inside a JavaScript single-page app — hard to use in the moment, impossible to reuse programmatically.", 0),
        ("Why it matters: wrong-bin items contaminate recycling and compost, raise cost, and undercut the city's waste-reduction goals.", 0),
        ("Goal: make the SAME official guidance portable so one dataset drives many trustworthy experiences — search, app, sign, voice, classroom.", 0),
        ("In scope: Seattle item-level disposal guidance as open JSON.  Out of scope: making the decision for the user, non-Seattle jurisdictions, accounts/PII.", 0),
    ]),

    ("2 · Existing Structure & FAIR Assessment", "What's there today, and where it's weak", [
        ("Existing form: category-nested website content; disposal logic buried in free-text paragraphs; data served by two undocumented taxonomy endpoints.", 0),
        ("Access: human navigation through a SPA; no stable IDs, no schema, no contract for reuse.", 0),
        ("FAIR assessment of the existing structure:", 0),
        ("Strong on Authority — SPU is the domain standard (R1.3 = High).", 1),
        ("Weak on Findable / Interoperable / Reusable — no portable IDs, no shared schema, no machine-readable index (F1, F4, I1 = Low).", 1),
        ("Those weak FAIR facets are exactly what the transformation targets.", 0),
    ]),

    ("3 · How We Decided to Improve It", "Deficiencies → transformation decisions", [
        ("Promote a stable item_id to a first-class key → dedup, joins, updates (fixes F1).", 0),
        ("Flatten category-nested pages into one uniform item array → searchable index (fixes F4, I1).", 0),
        ("Replace free-text paragraphs with structured disposal_streams + per-stream conditions; preserve ALL valid streams (74% of items are multi-stream) instead of forcing one answer.", 0),
        ("Attach per-item provenance (source_reference, city_jurisdiction) and a versioned metadata header (fixes F3, A2, R1.2).", 0),
        ("Validate the decision with a second city (Washington, DC): same schema, streams inferred from ReCollect templates — proves interoperability.", 0),
    ]),

    ("4 · What the New Structure Is", "Portable, item-keyed JSON — differs on all four dimensions", [
        ("One JSON object per item: item_id, item_name, synonyms, disposal_streams[], disposal_conditions{}, explanation, category_path[], source_reference, city_jurisdiction.", 0),
        ("310 Seattle items (+468 DC), schema v1.0, with a metadata header and the SPU category tree.", 0),
        ("Substantial difference (rubric needs ≥2 — this changes all 4):", 0),
        ("Information: normalized records w/ explicit multi-stream logic.  Structure: flat array vs. nested pages.", 1),
        ("Format: portable versioned JSON vs. HTML/JS payload.  Access: notebook + snapshot file + 9-endpoint REST API vs. SPA navigation.", 1),
        ("Meets every requirement (F/A/I/R) with no scope creep; SPU quirks preserved, not silently rewritten.", 0),
    ]),

    ("5 · How Quality Is Identified & Addressed", "Targets, measured results, and remediation", [
        ("Quality (measured on the snapshot): 310/310 unique IDs, 0 duplicates, 0 unknown streams, 0 stream/condition mismatches, 0.39 MB (< 2 MB target).", 0),
        ("Known gaps documented with remediation: 1 item with no stream (flagged, not guessed); 73 empty SPU summaries (conditions still present; derive fallback).", 0),
        ("Performance (measured vs. target): single lookup 1.1 ms (< 50), list 1.5 ms (< 200), search 1.4 ms (< 250), startup 0.26 s; 50/50 requests OK, zero 500s.", 0),
        ("Identified via TESTPLAN.md (functional + performance + quality tests, alarms & actions); reproducible commands in docs/05.", 0),
        ("Accessible & published: public GitHub repo with README, notebook, snapshot, API, and supporting docs for all five graded areas.", 0),
    ]),
]


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # green band
    left, top, width, height = Inches(0), Inches(2.2), prs.slide_width, Inches(1.6)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), prs.slide_width - Inches(1.2), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = title
    run.font.size = Pt(40); run.font.bold = True; run.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(16); r2.font.color.rgb = GREY
    return slide


def add_bullet_slide(prs, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(1.1))
    tf = tbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = GREEN
    if subtitle:
        sp = tf.add_paragraph()
        sr = sp.add_run(); sr.text = subtitle
        sr.font.size = Pt(15); sr.font.italic = True; sr.font.color.rgb = GREY
    # Body
    bbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), prs.slide_width - Inches(1.4), Inches(5.2))
    bf = bbox.text_frame; bf.word_wrap = True
    first = True
    for text, level in bullets:
        p = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        p.level = level
        bullet = "•  " if level == 0 else "–  "
        run = p.add_run(); run.text = bullet + text
        run.font.size = Pt(16 if level == 0 else 14)
        run.font.color.rgb = DARK if level == 0 else GREY
        p.space_after = Pt(8)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for title, subtitle, bullets in SLIDES:
        if bullets is None:
            add_title_slide(prs, title, subtitle)
        else:
            add_bullet_slide(prs, title, subtitle, bullets)
    prs.save(OUT)
    print("wrote", OUT, "with", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    main()
